from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colors ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x0e, 0x0e, 0x0e)   # dark bg
PRIMARY = RGBColor(0x00, 0x62, 0x33)   # deep green
ACCENT  = RGBColor(0x00, 0xa3, 0x55)   # bright green
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xa5, 0xd6, 0xa7)
MUTED   = RGBColor(0x6a, 0xad, 0x6a)
CARD    = RGBColor(0x14, 0x14, 0x14)
BORDER  = RGBColor(0x2a, 0x4a, 0x2a)
YELLOW  = RGBColor(0xff, 0xd5, 0x4f)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, l, t, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.RIGHT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def accent_line(slide, l, t, w=1.5):
    r = rect(slide, l, t, w, 0.05, ACCENT)
    return r

def card(slide, l, t, w, h):
    r = rect(slide, l, t, w, h, CARD)
    r.line.color.rgb = BORDER
    r.line.width = Pt(1)
    return r

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
bg(s1)
rect(s1, 0, 0, 13.33, 7.5, BG)
# left green bar
rect(s1, 0, 0, 0.12, 7.5, PRIMARY)
# glow circle (top-right)
glow = s1.shapes.add_shape(1, Inches(9), Inches(-1), Inches(5), Inches(5))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x00,0x30,0x18)
glow.line.fill.background()

txb(s1, "منصة سند", 1, 1.5, 11, 1.5, size=54, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
txb(s1, "المنصة الجزائرية الإعلامية الشاملة", 1, 3.0, 11, 0.8, size=24, color=LIGHT, align=PP_ALIGN.RIGHT)
accent_line(s1, 10.5, 3.9, 1.7)
txb(s1, "sanadz.media", 1, 4.2, 11, 0.6, size=16, color=MUTED, align=PP_ALIGN.RIGHT, italic=True)
txb(s1, "2026", 1, 6.5, 11, 0.6, size=13, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bg(s2)
rect(s2, 0, 0, 0.12, 7.5, PRIMARY)

txb(s2, "المشكلة", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s2, "قطاع الإعلام في الجزائر يفتقر إلى منصة رقمية موحّدة", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s2, 10.5, 1.75)

problems = [
    ("🔍", "تشتّت الفرص",       "فرص التوظيف والتدريب مبعثرة على عشرات المواقع والصفحات"),
    ("📚", "غياب التدريب",      "لا توجد منصة متخصصة تجمع الدورات الإعلامية الجزائرية"),
    ("🤝", "ضعف التواصل",       "الإعلاميون يصعب عليهم التواصل مع بعضهم والتعاون"),
    ("🛒", "سوق المعدات",        "لا يوجد سوق رقمي موثوق لمعدات الإعلام في الجزائر"),
]
for i, (icon, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    l = 6.8 - col * 6.2
    t = 2.2 + row * 2.1
    card(s2, l, t, 5.8, 1.85)
    txb(s2, icon + "  " + title, l+0.2, t+0.15, 5.4, 0.5, size=15, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    txb(s2, desc, l+0.2, t+0.55, 5.4, 1.0, size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bg(s3)
rect(s3, 0, 0, 0.12, 7.5, PRIMARY)

txb(s3, "الحل", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s3, "سند — المنصة الجزائرية الإعلامية الشاملة", 0.4, 0.9, 12, 0.8, size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s3, 10.5, 1.75)
txb(s3,
    "منصة رقمية واحدة تجمع الإعلاميين الجزائريين وتوفر لهم كل ما يحتاجونه:\n"
    "التوظيف · التدريب · المعدات · المنشطون · المسابقات · الأخبار",
    0.4, 2.0, 12, 1.2, size=16, color=LIGHT, align=PP_ALIGN.RIGHT)

features = [
    "💼 توظيف إعلامي",
    "🎓 دورات تدريبية",
    "📦 سوق المعدات",
    "🎙️ المنشطون",
    "🏆 مسابقات إعلامية",
    "📺 دليل القنوات",
]
for i, feat in enumerate(features):
    col = i % 3
    row = i // 3
    l = 12.6 - col * 4.1
    t = 3.5 + row * 1.6
    c = card(s3, l-3.8, t, 3.6, 1.3)
    txb(s3, feat, l-3.6, t+0.35, 3.2, 0.7, size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Features detail
# ════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
bg(s4)
rect(s4, 0, 0, 0.12, 7.5, PRIMARY)

txb(s4, "خدمات المنصة", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s4, "8 خدمات متكاملة في منصة واحدة", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s4, 10.5, 1.75)

services = [
    ("التوظيف الإعلامي",    "فرص عمل من قنوات ومؤسسات إعلامية في كل الجزائر"),
    ("الدورات التدريبية",   "برامج تدريبية من مدربين معتمدين في الصحافة والإنتاج"),
    ("سوق المعدات",          "بيع وشراء معدات التصوير والصوت والبث"),
    ("المنشطون",             "قاعدة بيانات لمنشطي البرامج والفعاليات"),
    ("المسابقات",            "مسابقات إعلامية محلية ودولية للطلاب والمحترفين"),
    ("دليل القنوات",         "دليل شامل للقنوات الجزائرية التلفزيونية والإذاعية"),
    ("مذكرات التخرج",        "بنك للأبحاث ومذكرات التخرج في تخصصات الإعلام"),
    ("الأخبار",              "آخر أخبار قطاع الإعلام في الجزائر"),
]
for i, (title, desc) in enumerate(services):
    col = i % 2
    row = i // 2
    l = 12.6 - col * 6.2
    t = 2.1 + row * 1.25
    card(s4, l-5.8, t, 5.6, 1.1)
    txb(s4, "● " + title, l-5.6, t+0.1, 5.2, 0.4, size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    txb(s4, desc, l-5.6, t+0.5, 5.2, 0.5, size=11, color=LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Target audience
# ════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
bg(s5)
rect(s5, 0, 0, 0.12, 7.5, PRIMARY)

txb(s5, "الجمهور المستهدف", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s5, "من يستفيد من سند؟", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s5, 10.5, 1.75)

audiences = [
    ("👨‍💻", "الإعلاميون المحترفون",
     "صحفيون، مصورون، مخرجون، مونتيرون\nيبحثون عن فرص عمل وتطوير مهني"),
    ("🎓", "طلاب الإعلام",
     "طلاب الجامعات الجزائرية في تخصصات\nالإعلام والاتصال والصحافة"),
    ("📡", "المؤسسات الإعلامية",
     "قنوات تلفزيونية، إذاعات، وكالات أنباء\nتبحث عن كفاءات إعلامية"),
    ("🏫", "مراكز التدريب",
     "مدربون ومراكز تقدم دورات تدريبية\nفي مجال الإعلام والإنتاج"),
]
for i, (icon, title, desc) in enumerate(audiences):
    l = 12.5 - i * 3.05
    card(s5, l-2.7, 2.2, 2.8, 4.5)
    txb(s5, icon, l-2.5, 2.4, 2.4, 0.7, size=28, align=PP_ALIGN.CENTER)
    txb(s5, title, l-2.5, 3.1, 2.4, 0.6, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s5, desc, l-2.5, 3.75, 2.4, 2.5, size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Market
# ════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
bg(s6)
rect(s6, 0, 0, 0.12, 7.5, PRIMARY)

txb(s6, "السوق المستهدف", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s6, "قطاع الإعلام الجزائري — فرصة ضخمة", 0.4, 0.9, 12, 0.8, size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s6, 10.5, 1.75)

stats = [
    ("69", "ولاية جزائرية", "تغطية كاملة للوطن"),
    ("+500", "وسيلة إعلامية", "قنوات، إذاعات، مواقع إخبارية"),
    ("+30,000", "طالب إعلام", "مسجّلون في الجامعات سنوياً"),
    ("+10,000", "إعلامي محترف", "ينشطون في القطاع"),
]
for i, (num, label, sub) in enumerate(stats):
    l = 12.5 - i * 3.05
    card(s6, l-2.7, 2.2, 2.8, 2.8)
    txb(s6, num, l-2.7, 2.4, 2.8, 1.0, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s6, label, l-2.7, 3.35, 2.8, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s6, sub, l-2.7, 3.85, 2.8, 0.7, size=10, color=MUTED, align=PP_ALIGN.CENTER)

txb(s6,
    "الجزائر تمتلك أكثر من 60 قناة تلفزيونية وأكثر من 100 إذاعة — كلها تحتاج كوادر بشرية مؤهلة",
    0.5, 5.4, 12.3, 0.8, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Business Model
# ════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
bg(s7)
rect(s7, 0, 0, 0.12, 7.5, PRIMARY)

txb(s7, "نموذج الأعمال", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s7, "مصادر الإيرادات", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s7, 10.5, 1.75)

models = [
    ("💰", "اشتراك المتاجر",
     "المتاجر المحترفة تدفع اشتراكاً شهرياً\nللظهور في سوق المعدات"),
    ("🎓", "عمولة الدورات",
     "نسبة مئوية من رسوم الدورات\nالتي تُباع عبر المنصة"),
    ("📢", "الإعلانات",
     "إعلانات المؤسسات الإعلامية\nوالشركات الراغبة في التوظيف"),
    ("⭐", "الإبراز المدفوع",
     "خدمة الإبراز للمتاجر والمدربين\nوالمنشطين في نتائج البحث"),
]
for i, (icon, title, desc) in enumerate(models):
    col = i % 2
    row = i // 2
    l = 12.5 - col * 6.1
    t = 2.2 + row * 2.3
    card(s7, l-5.8, t, 5.6, 2.0)
    txb(s7, icon + "  " + title, l-5.6, t+0.2, 5.2, 0.5, size=15, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)
    txb(s7, desc, l-5.6, t+0.7, 5.2, 1.0, size=12, color=LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
bg(s8)
rect(s8, 0, 0, 0.12, 7.5, PRIMARY)

txb(s8, "التقنية المستخدمة", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s8, "منصة حديثة وقابلة للتوسع", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s8, 10.5, 1.75)

techs = [
    ("React + TypeScript",  "واجهة مستخدم سريعة وتفاعلية"),
    ("Firebase Firestore",  "قاعدة بيانات سحابية آنية"),
    ("Firebase Auth",       "نظام تسجيل دخول آمن"),
    ("Cloudinary",          "تخزين وتحسين الصور"),
    ("PWA",                 "تطبيق ويب تقدمي — يعمل بدون إنترنت"),
    ("Vercel",              "نشر سريع مع CDN عالمي"),
    ("FCM",                 "نظام الإشعارات الفورية"),
    ("TailwindCSS",         "تصميم متجاوب لكل الأجهزة"),
]
for i, (tech, desc) in enumerate(techs):
    col = i % 2
    row = i // 2
    l = 12.5 - col * 6.1
    t = 2.2 + row * 1.2
    card(s8, l-5.8, t, 5.6, 1.0)
    txb(s8, tech, l-5.6, t+0.08, 2.5, 0.45, size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    txb(s8, desc, l-3.0, t+0.08, 2.7, 0.45, size=12, color=LIGHT, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Traction
# ════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
bg(s9)
rect(s9, 0, 0, 0.12, 7.5, PRIMARY)

txb(s9, "ما تم إنجازه", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s9, "المنصة جاهزة وتعمل", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s9, 10.5, 1.75)

done = [
    "✅  المنصة مطلوقة رسمياً على sanadz.media",
    "✅  8 خدمات رقمية متكاملة تعمل بشكل كامل",
    "✅  نظام تسجيل للإعلاميين والمتاجر ومراكز التدريب",
    "✅  نظام إشعارات فورية (Push Notifications)",
    "✅  لوحة تحكم إدارية متكاملة للمدير",
    "✅  لوحة تحكم لكل مستخدم لإدارة ملفه ومنتجاته ودوراته",
    "✅  تطبيق ويب تقدمي (PWA) — يدعم التثبيت على الجوال",
    "✅  تصميم RTL كامل باللغة العربية",
    "✅  نظام موافقة إدارية على كل المحتوى",
]
for i, item in enumerate(done):
    t = 2.1 + i * 0.55
    txb(s9, item, 0.4, t, 12.3, 0.5, size=13, color=LIGHT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Roadmap
# ════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
bg(s10)
rect(s10, 0, 0, 0.12, 7.5, PRIMARY)

txb(s10, "خارطة الطريق", 0.4, 0.3, 12, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
txb(s10, "الخطوات القادمة", 0.4, 0.9, 12, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
accent_line(s10, 10.5, 1.75)

phases = [
    ("المرحلة 1\nالمنجزة ✅", "Q2 2026",
     "إطلاق المنصة\nالخدمات الأساسية الثمانية\nنظام التسجيل والإدارة"),
    ("المرحلة 2\nقيد التطوير", "Q3 2026",
     "تطبيق موبايل (Android/iOS)\nنظام الدفع الإلكتروني\nتوسيع قاعدة المستخدمين"),
    ("المرحلة 3\nمستقبلية", "Q4 2026",
     "شراكات مع مؤسسات إعلامية\nنظام التوصيات بالذكاء الاصطناعي\nالتوسع لدول المغرب العربي"),
]
colors_ph = [ACCENT, YELLOW, MUTED]
for i, (phase, period, details) in enumerate(phases):
    l = 12.2 - i * 4.1
    card(s10, l-3.8, 2.2, 3.6, 4.5)
    rect(s10, l-3.8, 2.2, 3.6, 0.08, colors_ph[i])
    txb(s10, phase, l-3.6, 2.3, 3.2, 0.8, size=14, bold=True, color=colors_ph[i], align=PP_ALIGN.CENTER)
    txb(s10, period, l-3.6, 3.1, 3.2, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s10, details, l-3.6, 3.6, 3.2, 2.5, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Ask / CTA
# ════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
bg(s11)
rect(s11, 0, 0, 0.12, 7.5, PRIMARY)

# big accent glow
glow2 = s11.shapes.add_shape(1, Inches(4), Inches(1.5), Inches(5.5), Inches(5.5))
glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x00,0x25,0x12)
glow2.line.fill.background()

txb(s11, "ما نحتاجه", 0.4, 0.5, 12.3, 0.7, size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
txb(s11, "انضم إلى رحلة سند", 0.4, 1.1, 12.3, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
accent_line(s11, 5.9, 2.1, 1.5)

needs = [
    ("🏛️", "دعم منصة الشركات الناشئة", "الاستفادة من برامج الدعم والتمويل الحكومي"),
    ("🤝", "شراكات مؤسسية",           "التعاون مع وزارة الاتصال والمؤسسات الإعلامية"),
    ("📣", "انتشار وتسويق",             "الوصول إلى أكبر عدد من الإعلاميين الجزائريين"),
]
for i, (icon, title, desc) in enumerate(needs):
    l_positions = [12.4, 8.1, 3.8]
    l = l_positions[i]
    card(s11, l-3.5, 2.8, 3.3, 3.0)
    txb(s11, icon, l-3.3, 3.0, 2.9, 0.7, size=26, align=PP_ALIGN.CENTER)
    txb(s11, title, l-3.3, 3.7, 2.9, 0.6, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s11, desc, l-3.3, 4.35, 2.9, 1.1, size=11, color=LIGHT, align=PP_ALIGN.CENTER)

txb(s11, "sanadz.media  |  officialelvaz@gmail.com", 0.4, 6.5, 12.3, 0.6, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
prs.save("/home/user/sanad/Sanad_Pitch_Deck.pptx")
print("Done!")
